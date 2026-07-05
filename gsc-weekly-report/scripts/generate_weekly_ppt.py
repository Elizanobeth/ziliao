#!/usr/bin/env python3
"""Fill a GSC weekly report PPTX template from lossless organized JSON."""

from __future__ import annotations

import argparse
import copy
import datetime as dt
import json
import re
import sys
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError as exc:  # pragma: no cover
    raise SystemExit("Missing dependency: python-pptx. Install python-pptx or use the Codex bundled Python runtime.") from exc


PLACEHOLDER_RE = re.compile(r"\{\{[A-Z0-9_]+\}\}")

CATEGORIES = {
    "main_progress": {
        "title": "主要进展",
        "placeholder": "{{MAIN_PROGRESS}}",
        "empty_text": "暂无主要进展。",
    },
    "issues": {
        "title": "存在问题",
        "placeholder": "{{ISSUES}}",
        "empty_text": "暂无明显问题。",
    },
    "help_needed": {
        "title": "需要帮助",
        "placeholder": "{{HELP_NEEDED}}",
        "empty_text": "暂无需额外协调事项。",
    },
    "next_plan": {
        "title": "下周计划",
        "placeholder": "{{NEXT_PLAN}}",
        "empty_text": "暂无明确下周计划。",
    },
}


def ensure_list(value: Any) -> list[str]:
    if value is None:
        return []
    if isinstance(value, list):
        return [item_to_text(item) for item in value if item_to_text(item)]
    text = str(value).strip()
    return [text] if text else []


def item_to_text(item: Any) -> str:
    if item is None:
        return ""
    if isinstance(item, dict):
        for key in ("text", "content", "summary", "item"):
            if item.get(key):
                return str(item[key]).strip()
        return "；".join(f"{key}：{value}" for key, value in item.items() if value)
    return str(item).strip()


def clean_item_prefix(text: str) -> str:
    return re.sub(r"^[•·\-\d\.、）\)\s]+", "", text).strip()


def normalized_category_items(summary: dict[str, Any], category: str) -> list[tuple[int, str]]:
    raw_items = ensure_list(summary.get(category))
    if not raw_items:
        raw_items = [CATEGORIES[category]["empty_text"]]

    normalized: list[tuple[int, str]] = []
    item_number = 1
    for raw_item in raw_items:
        normalized.append((item_number, clean_item_prefix(raw_item)))
        item_number += 1
    return normalized


def empty_page() -> dict[str, list[tuple[int, str]]]:
    return {category: [] for category in CATEGORIES}


def page_has_content(page: dict[str, list[tuple[int, str]]]) -> bool:
    return any(page[category] for category in CATEGORIES)


def category_char_count(items: list[tuple[int, str]]) -> int:
    return sum(len(text) for _, text in items)


def paginate_summary(
    summary: dict[str, Any],
    max_items_per_category: int,
    max_chars_per_category: int,
) -> list[dict[str, list[tuple[int, str]]]]:
    category_items = {
        category: normalized_category_items(summary, category)
        for category in CATEGORIES
    }
    pages: list[dict[str, list[tuple[int, str]]]] = []
    page = empty_page()

    for category, items in category_items.items():
        position = 0
        while position < len(items):
            if len(page[category]) >= max_items_per_category:
                pages.append(page)
                page = empty_page()
                continue

            current_chars = category_char_count(page[category])
            item_number, text = items[position]
            item_chars = len(text)
            would_exceed_chars = current_chars > 0 and current_chars + item_chars > max_chars_per_category
            if would_exceed_chars:
                pages.append(page)
                page = empty_page()
                continue

            page[category].append((item_number, text))
            position += 1

            if position < len(items):
                current_chars = category_char_count(page[category])
                if (
                    len(page[category]) >= max_items_per_category
                    or current_chars >= max_chars_per_category
                ):
                    pages.append(page)
                    page = empty_page()

    if page_has_content(page):
        pages.append(page)

    return pages or [empty_page()]


def bullet_text(items: list[tuple[int, str]]) -> str:
    if not items:
        return ""
    rendered = []
    for item_number, item in items:
        clean = clean_item_prefix(item)
        rendered.append(f"{item_number}）{clean}")
    return "\n".join(rendered)


def build_replacements(
    summary: dict[str, Any],
    page_items: dict[str, list[tuple[int, str]]],
    page_number: int,
    page_count: int,
) -> dict[str, str]:
    week = str(summary.get("week") or "").strip()
    title = str(summary.get("title") or f"{week} 周报").strip()
    generated_date = str(summary.get("generated_date") or dt.date.today().isoformat())
    scope = str(summary.get("scope") or "全部section").strip()
    replacements = {
        "{{REPORT_WEEK}}": week,
        "{{REPORT_TITLE}}": title,
        "{{SCOPE}}": scope,
        "{{GENERATED_DATE}}": generated_date,
    }
    for category, config in CATEGORIES.items():
        replacements[config["placeholder"]] = bullet_text(page_items[category])
    return replacements


def replace_text(text: str, replacements: dict[str, str]) -> str:
    for placeholder, value in replacements.items():
        text = text.replace(placeholder, value)
    return text


def font_snapshot(text_frame: Any) -> dict[str, Any]:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            font = run.font
            color = None
            try:
                color = font.color.rgb
            except AttributeError:
                color = None
            return {
                "name": font.name,
                "size": font.size,
                "bold": font.bold,
                "italic": font.italic,
                "underline": font.underline,
                "color": color,
            }
    return {}


def apply_font(run: Any, snapshot: dict[str, Any]) -> None:
    font = run.font
    if snapshot.get("name"):
        font.name = snapshot["name"]
    if snapshot.get("size"):
        font.size = snapshot["size"]
    if snapshot.get("bold") is not None:
        font.bold = snapshot["bold"]
    if snapshot.get("italic") is not None:
        font.italic = snapshot["italic"]
    if snapshot.get("underline") is not None:
        font.underline = snapshot["underline"]
    if snapshot.get("color") is not None:
        font.color.rgb = snapshot["color"]


def set_text_preserving_style(text_frame: Any, text: str) -> None:
    style = font_snapshot(text_frame)
    alignment = text_frame.paragraphs[0].alignment if text_frame.paragraphs else None
    text_frame.clear()
    lines = text.splitlines() or [""]
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = alignment
        run = paragraph.add_run()
        run.text = line
        apply_font(run, style)


def iter_shapes(shapes: Any):
    for shape in shapes:
        yield shape
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def direct_shapes(slide: Any) -> list[Any]:
    return list(slide.shapes)


def shape_text(shape: Any) -> str:
    texts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        texts.append(shape.text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                texts.append(cell.text)
    return "\n".join(texts)


def shape_bounds(shape: Any) -> tuple[int, int, int, int]:
    left = int(shape.left)
    top = int(shape.top)
    return left, top, left + int(shape.width), top + int(shape.height)


def union_bounds(bounds: list[tuple[int, int, int, int]]) -> tuple[int, int, int, int]:
    return (
        min(bound[0] for bound in bounds),
        min(bound[1] for bound in bounds),
        max(bound[2] for bound in bounds),
        max(bound[3] for bound in bounds),
    )


def expand_bounds(bounds: tuple[int, int, int, int], margin: int) -> tuple[int, int, int, int]:
    return bounds[0] - margin, bounds[1] - margin, bounds[2] + margin, bounds[3] + margin


def bounds_intersect(a: tuple[int, int, int, int], b: tuple[int, int, int, int]) -> bool:
    return a[0] < b[2] and a[2] > b[0] and a[1] < b[3] and a[3] > b[1]


def collect_category_groups(slide: Any) -> dict[str, dict[str, Any]]:
    shapes = direct_shapes(slide)
    groups: dict[str, dict[str, Any]] = {}
    for category, config in CATEGORIES.items():
        anchors = [
            shape
            for shape in shapes
            if config["placeholder"] in shape_text(shape) or config["title"] in shape_text(shape)
        ]
        if not anchors:
            groups[category] = {"shapes": [], "bounds": None}
            continue

        anchor_bounds = union_bounds([shape_bounds(shape) for shape in anchors])
        search_bounds = expand_bounds(anchor_bounds, 120000)
        group_shapes = [
            shape
            for shape in shapes
            if bounds_intersect(shape_bounds(shape), search_bounds)
        ]
        group_bounds = union_bounds([shape_bounds(shape) for shape in group_shapes])
        groups[category] = {"shapes": group_shapes, "bounds": group_bounds}
    return groups


def remove_shapes(slide: Any, shapes: list[Any]) -> None:
    seen: set[int] = set()
    for shape in shapes:
        key = id(shape.element)
        if key in seen:
            continue
        seen.add(key)
        try:
            slide.shapes._spTree.remove(shape.element)  # noqa: SLF001
        except ValueError:
            pass


def move_shapes(shapes: list[Any], dx: int, dy: int) -> None:
    for shape in shapes:
        shape.left = int(shape.left) + dx
        shape.top = int(shape.top) + dy


def adjust_category_layout(slide: Any, page_items: dict[str, list[tuple[int, str]]]) -> None:
    groups = collect_category_groups(slide)
    slot_bounds = [
        groups[category]["bounds"]
        for category in CATEGORIES
        if groups[category]["bounds"] is not None
    ]
    if not slot_bounds:
        return

    empty_categories = [category for category in CATEGORIES if not page_items[category]]
    for category in empty_categories:
        remove_shapes(slide, groups[category]["shapes"])

    visible_categories = [category for category in CATEGORIES if page_items[category]]
    for target_index, category in enumerate(visible_categories):
        group = groups[category]
        if not group["shapes"] or group["bounds"] is None or target_index >= len(slot_bounds):
            continue
        target_bounds = slot_bounds[target_index]
        current_bounds = group["bounds"]
        dx = target_bounds[0] - current_bounds[0]
        dy = target_bounds[1] - current_bounds[1]
        if dx or dy:
            move_shapes(group["shapes"], dx, dy)


def replace_in_text_frame(shape: Any, replacements: dict[str, str]) -> int:
    if not getattr(shape, "has_text_frame", False):
        return 0
    original = shape.text
    updated = replace_text(original, replacements)
    if updated == original:
        return 0
    set_text_preserving_style(shape.text_frame, updated)
    return 1


def replace_in_table(shape: Any, replacements: dict[str, str]) -> int:
    if not getattr(shape, "has_table", False):
        return 0
    changed = 0
    for row in shape.table.rows:
        for cell in row.cells:
            original = cell.text
            updated = replace_text(original, replacements)
            if updated != original:
                cell.text = updated
                changed += 1
    return changed


def unresolved_placeholders(prs: Presentation) -> list[str]:
    unresolved: set[str] = set()
    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            texts: list[str] = []
            if getattr(shape, "has_text_frame", False):
                texts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text)
            for text in texts:
                unresolved.update(PLACEHOLDER_RE.findall(text))
    return sorted(unresolved)


def delete_slide(prs: Presentation, slide_index: int) -> None:
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public delete API.
    slides = list(slide_id_list)
    slide_id_list.remove(slides[slide_index])


def duplicate_slide(prs: Presentation, source_slide: Any) -> Any:
    copied_slide = prs.slides.add_slide(source_slide.slide_layout)
    for shape in list(copied_slide.shapes):
        copied_slide.shapes._spTree.remove(shape.element)  # noqa: SLF001
    for shape in source_slide.shapes:
        copied_element = copy.deepcopy(shape.element)
        copied_slide.shapes._spTree.insert_element_before(copied_element, "p:extLst")  # noqa: SLF001
    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        try:
            new_rid = copied_slide.part.rels._add_relationship(  # noqa: SLF001
                rel.reltype,
                rel._target,  # noqa: SLF001
                getattr(rel, "is_external", False),
            )
            for element in copied_slide.element.iter():
                for attr_name, attr_value in list(element.attrib.items()):
                    if attr_value == rel.rId:
                        element.attrib[attr_name] = new_rid
        except ValueError:
            pass
    return copied_slide


def prepare_paginated_slides(prs: Presentation, page_count: int) -> None:
    if len(prs.slides) == 0:
        raise ValueError("Template PPTX must contain at least one slide.")
    while len(prs.slides) > 1:
        delete_slide(prs, len(prs.slides) - 1)
    template_slide = prs.slides[0]
    while len(prs.slides) < page_count:
        duplicate_slide(prs, template_slide)


def fill_template(
    template_path: Path,
    summary_path: Path,
    output_path: Path,
    max_items_per_category: int,
    max_chars_per_category: int,
) -> dict[str, Any]:
    summary = json.loads(summary_path.read_text(encoding="utf-8"))
    pages = paginate_summary(summary, max_items_per_category, max_chars_per_category)
    prs = Presentation(str(template_path))
    prepare_paginated_slides(prs, len(pages))

    changed = 0
    for page_number, (slide, page_items) in enumerate(zip(prs.slides, pages), start=1):
        adjust_category_layout(slide, page_items)
        replacements = build_replacements(summary, page_items, page_number, len(pages))
        for shape in iter_shapes(slide.shapes):
            changed += replace_in_text_frame(shape, replacements)
            changed += replace_in_table(shape, replacements)

    unresolved = unresolved_placeholders(prs)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return {
        "output": str(output_path),
        "slide_count": len(pages),
        "changed_shapes_or_cells": changed,
        "unresolved_placeholders": unresolved,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate a GSC weekly report PPTX from lossless organized JSON.")
    parser.add_argument("--summary", required=True, help="Path to weekly_report_summary.json.")
    parser.add_argument("--template", required=True, help="Path to PPTX template with placeholders.")
    parser.add_argument("--output", required=True, help="Output PPTX path.")
    parser.add_argument(
        "--max-items-per-category-per-slide",
        type=int,
        default=3,
        help="Maximum organized items to place in each category on one slide.",
    )
    parser.add_argument(
        "--max-chars-per-category-per-slide",
        type=int,
        default=220,
        help="Approximate maximum Chinese characters to place in each category on one slide.",
    )
    args = parser.parse_args()

    summary_path = Path(args.summary).expanduser().resolve()
    template_path = Path(args.template).expanduser().resolve()
    output_path = Path(args.output).expanduser().resolve()

    if not summary_path.exists():
        raise SystemExit(f"Summary JSON does not exist: {summary_path}")
    if not template_path.exists():
        raise SystemExit(f"Template PPTX does not exist: {template_path}")

    if args.max_items_per_category_per_slide < 1:
        raise SystemExit("--max-items-per-category-per-slide must be >= 1")
    if args.max_chars_per_category_per_slide < 40:
        raise SystemExit("--max-chars-per-category-per-slide must be >= 40")

    result = fill_template(
        template_path,
        summary_path,
        output_path,
        args.max_items_per_category_per_slide,
        args.max_chars_per_category_per_slide,
    )
    print(json.dumps(result, ensure_ascii=False))
    return 0 if not result["unresolved_placeholders"] else 2


if __name__ == "__main__":
    sys.exit(main())
